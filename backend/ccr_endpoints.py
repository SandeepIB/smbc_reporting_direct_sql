from fastapi import HTTPException, UploadFile
from fastapi.responses import FileResponse
from typing import List, Dict
import os
import shutil
import tempfile
from datetime import datetime
from openai import OpenAI
import base64
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re
import io

# Global variables
analysis_results = {}
temp_dir = None
crop_config = {"rows": 2, "cols": 3, "enabled": True}
selected_template = "SMBC.pptx"  # Default template

# OpenAI client
api_key = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=api_key) if api_key else None

# Template management functions
def get_available_templates() -> List[Dict]:
    """Get list of available PPT templates with metadata"""
    templates_dir = os.path.join(os.path.dirname(__file__), 'templates')
    print(f"Looking for templates in: {templates_dir}")
    print(f"Templates directory exists: {os.path.exists(templates_dir)}")
    if os.path.exists(templates_dir):
        print(f"Files in templates directory: {os.listdir(templates_dir)}")
    templates = []
    
    template_info = {
        "SMBC.pptx": {
            "name": "SMBC Executive Summary",
            "description": "Official SMBC template for executive reporting",
            "active": True
        },
        "CCR Executive Summary.pptx": {
            "name": "CCR Executive Summary", 
            "description": "Sample template for CCR executive reporting",
            "active": False
        },
        "Individual CCR Summary.pptx": {
            "name": "Individual CCR Summary",
            "description": "InfoBeans demo template for individual CCR analysis", 
            "active": False
        }
    }
    
    if not os.path.exists(templates_dir):
        print(f"Templates directory not found: {templates_dir}")
        return []
    
    for filename in os.listdir(templates_dir):
        if filename.endswith('.pptx'):
            template_path = os.path.join(templates_dir, filename)
            info = template_info.get(filename, {
                "name": filename.replace('.pptx', ''),
                "description": "Custom template",
                "active": False
            })
            
            # Generate preview
            preview = generate_template_preview(template_path)
            
            templates.append({
                "filename": filename,
                "name": info["name"],
                "description": info["description"],
                "active": info["active"],
                "path": template_path,
                "preview": preview
            })
    
    return templates

def generate_template_preview(template_path: str) -> str:
    """Get thumbnail image from templates folder"""
    try:
        # Get template filename without extension
        template_name = os.path.splitext(os.path.basename(template_path))[0]
        
        # Look for thumbnail image in templates folder
        templates_dir = os.path.dirname(template_path)
        
        # Try different image extensions
        for ext in ['.png', '.jpg', '.jpeg']:
            thumb_path = os.path.join(templates_dir, f"{template_name}{ext}")
            if os.path.exists(thumb_path):
                with open(thumb_path, 'rb') as f:
                    img_data = base64.b64encode(f.read()).decode('utf-8')
                return f"data:image/{ext[1:]};base64,{img_data}"
        
        return None
        
    except Exception as e:
        print(f"Error loading thumbnail for {template_path}: {e}")
        return None

# Image cropping functions
def crop_image_to_blocks(image_path, rows=2, cols=3, output_dir=None):
    if output_dir is None:
        output_dir = os.path.join(os.path.dirname(image_path), 'blocks')
    
    img = Image.open(image_path)
    img_width, img_height = img.size
    
    block_width = img_width // cols
    block_height = img_height // rows
    
    os.makedirs(output_dir, exist_ok=True)
    
    cropped_files = []
    for row in range(rows):
        for col in range(cols):
            left = col * block_width
            upper = row * block_height
            right = min(left + block_width, img_width)
            lower = min(upper + block_height, img_height)
            
            block = img.crop((left, upper, right, lower))
            
            if row == rows - 1 and col == cols - 1:
                block_name = 'key_risk_indicators.jpg'
            else:
                block_name = f'block_{row}_{col}.jpg'
            
            block_path = os.path.join(output_dir, block_name)
            block.save(block_path)
            cropped_files.append(block_path)
    
    return cropped_files

def process_uploaded_images(temp_dir, rows=2, cols=3):
    all_cropped_files = []
    
    for filename in os.listdir(temp_dir):
        if filename.lower().endswith(('.png', '.jpg', '.jpeg')):
            image_path = os.path.join(temp_dir, filename)
            
            base_name = os.path.splitext(filename)[0]
            blocks_dir = os.path.join(temp_dir, f'{base_name}_blocks')
            
            cropped_files = crop_image_to_blocks(image_path, rows, cols, blocks_dir)
            all_cropped_files.extend(cropped_files)
    
    return all_cropped_files

# Analysis functions
def analyze_image_with_gpt(image_path):
    if not client:
        raise Exception("OpenAI API key not configured")
        
    with open(image_path, "rb") as f:
        image_bytes = f.read()

    base64_image = base64.b64encode(image_bytes).decode("utf-8")
    user_context = "This Executive Summary report is intended for the risk management team of a multinational bank, providing a clear and structured view of Counterparty Credit Risk (CCR) exposure across Derivatives and Securities Financing Transactions (SFTs)."

    prompt_text = (
        f"You are an expert data analyst AI. A user uploads a graph image for your analysis.\n\n"
        f"User context: {user_context}\n\n"
        "Your task is to analyze the chart visually and generate two short but informative paragraphs.\n\n"
        "Respond in the following format using bold section headers:\n"
        "**Trend Detection:** Describe the overall patterns, changes, and deviations observed in the graph. "
        "Highlight trends over time, categories with significant movements, and any gender/location-based observations if applicable.\n\n"
        "**Recommendations:** Based on your observations, provide clear, actionable suggestions to the business. "
        "These should address potential risks, growth opportunities, or areas that require attention or improvement.\n\n"
        "Avoid repeating data values verbatim. Be concise, insightful, and practical."
    )

    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": prompt_text},
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "Analyze this image."},
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{base64_image}"}}
                ],
            }
        ],
        max_tokens=600,
        temperature=0.7
    )

    return response.choices[0].message.content

def analyze_all_graphs(folder):
    graph_insights = []

    # Get all image files
    all_files = []
    for root, dirs, files in os.walk(folder):
        for filename in files:
            if filename.lower().endswith((".png", ".jpg", ".jpeg")):
                all_files.append(os.path.join(root, filename))
    
    print(f"Found {len(all_files)} images to analyze")
    
    if not all_files:
        raise Exception("No image files found for analysis")
    
    for i, path in enumerate(sorted(all_files)):
        filename = os.path.basename(path)
        print(f"Analyzing {i+1}/{len(all_files)}: {filename}...")
        
        try:
            # Analyze image content
            result = analyze_image_with_gpt(path)
            
            # Parse trend and recommendation
            try:
                parts = result.split("**Recommendations:**")
                if len(parts) == 2:
                    trend = parts[0].replace("**Trend Detection:**", "").strip()
                    recommendation = parts[1].strip()
                else:
                    # Fallback parsing
                    trend_start = result.find("**Trend Detection:**")
                    rec_start = result.find("**Recommendations:**")
                    
                    if trend_start >= 0 and rec_start >= 0:
                        trend = result[trend_start + len("**Trend Detection:**"):rec_start].strip()
                        recommendation = result[rec_start + len("**Recommendations:**"):].strip()
                    else:
                        trend = result[:len(result)//2] if result else "Analysis completed."
                        recommendation = result[len(result)//2:] if result else "Review recommended."
            except Exception as parse_error:
                print(f"Parse error for {filename}: {parse_error}")
                trend = "Analysis shows data patterns requiring attention."
                recommendation = "Further review and monitoring recommended."
            
            # Get title
            try:
                title_prompt = "Extract only the main chart title from this image. Respond with just the title text, no other commentary."
                with open(path, "rb") as f:
                    image_bytes = f.read()

                base64_image = base64.b64encode(image_bytes).decode("utf-8")
                title_response = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": title_prompt},
                        {
                            "role": "user",
                            "content": [
                                {"type": "text", "text": "What is the title of this chart?"},
                                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{base64_image}"}}
                            ],
                        }
                    ],
                    max_tokens=50,
                    temperature=0.1,
                )
                title = title_response.choices[0].message.content.strip().strip('"').strip("'")
                if not title or len(title) < 3:
                    title = f"Chart Analysis: {filename}"
            except Exception as title_error:
                print(f"Title extraction error for {filename}: {title_error}")
                title = f"Chart Analysis: {filename}"
            
            graph_insights.append({
                'image_path': path,
                'trend': trend,
                'recommendation': recommendation,
                'title': title
            })
            
            print(f"✅ Successfully analyzed {filename}")
            
        except Exception as analysis_error:
            print(f"❌ Failed to analyze {filename}: {analysis_error}")
            # Add a fallback insight
            graph_insights.append({
                'image_path': path,
                'trend': "Chart analysis encountered technical difficulties.",
                'recommendation': "Manual review of this chart is recommended.",
                'title': f"Chart: {filename}"
            })
            
    all_trends = [insight['trend'] for insight in graph_insights]
    all_recommendations = [insight['recommendation'] for insight in graph_insights]

    consolidated_prompt = f"""
            You are a senior business analyst AI. Your task is to review multiple chart-level insights and provide a concise, insightful summary in two parts using the exact bold headers below:

            **Trend Detection:** Write a high-level executive summary paragraph that synthesizes and combines the overall patterns, changes, and key observations from the provided trends.

            **Recommendations:** Based on the trends above, provide a strategic recommendation paragraph that offers practical business guidance.

            Here are the individual insights to analyze:

            **Trend Detection Points:**  
            {chr(10).join(f"- {trend}" for trend in all_trends)}

            **Recommendation Points:**  
            {chr(10).join(f"- {rec}" for rec in all_recommendations)}
            """
    consolidated_response = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": consolidated_prompt}],
            max_tokens=500,
            temperature=0.7,
    )

    consolidated_summary = consolidated_response.choices[0].message.content
    
    consolidated_parts_trend = "Trend observation not found."
    consolidated_parts_recommendation = "Recommendation not found."
    try:
        consolidated_parts = consolidated_summary.split("**Recommendations:**")
        if len(consolidated_parts) == 2:
            consolidated_parts_trend = consolidated_parts[0].replace("**Trend Detection:**", "").strip()
            consolidated_parts_recommendation = consolidated_parts[1].strip()
    except IndexError:
        pass

    exec_summary = {
        "trend": consolidated_parts_trend,
        "recommendation": consolidated_parts_recommendation
    }

    return exec_summary, graph_insights

# PowerPoint generation functions
def add_bullets_with_bold(text_frame, text_block, font_size=10, color=RGBColor(255, 255, 255)):
    lines = text_block.strip().split('\n')

    for line in lines:
        line = line.strip().lstrip("-").strip()
        if not line:
            continue

        p = text_frame.add_paragraph()
        p.space_after = Pt(6)

        dash_run = p.add_run()
        dash_run.text = "- "
        dash_run.font.size = Pt(font_size)
        dash_run.font.color.rgb = color

        segments = re.split(r'(\*\*.*?\*\*)', line)
        for segment in segments:
            run = p.add_run()
            run.text = segment[2:-2] if segment.startswith("**") and segment.endswith("**") else segment
            run.font.bold = segment.startswith("**") and segment.endswith("**")
            run.font.size = Pt(font_size)
            run.font.color.rgb = color

def insert_slide_at(prs, slide, position=0):
    slide_id_list = prs.slides._sldIdLst
    slide_element = slide._element
    slide_id = slide_id_list[-1]
    slide_id_list.remove(slide_id)
    slide_id_list.insert(position, slide_id)

def add_trend_slide(prs, image_path, trend_texts, recommendations, trend_title, observation_title, recommendation_title, position=None):
    blank_layout = next((layout for layout in prs.slide_layouts if layout.name.lower() == "blank"), prs.slide_layouts[6])
    slide = prs.slides.add_slide(blank_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    background_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, slide_width, slide_height
    )
    background_shape.fill.solid()
    background_shape.fill.fore_color.rgb = RGBColor(255, 249, 237)
    background_shape.line.fill.background()
    slide.shapes._spTree.insert(2, background_shape._element)

    title_box = slide.shapes.add_textbox(Inches(0), Inches(0.3), Inches(5.5), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.clear()

    p = title_tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT

    run = p.add_run()
    run.text = trend_title
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(30, 30, 50)

    left_width = int(slide_width * 0.45)
    right_width = int(slide_width * 0.55)
    margin_top = Inches(1.2)

    slide.shapes.add_picture(image_path, Inches(0.5), margin_top, width=left_width - Inches(1))

    right_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left_width,
        0,
        slide_width - left_width,
        slide_height
    )
    fill = right_panel.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(102, 0, 0)
    right_panel.line.fill.background()

    trend_box = slide.shapes.add_textbox(
        left_width + Inches(0.2),
        Inches(0.5),
        right_width - Inches(0.4),
        Inches(3.2)
    )
    trend_tf = trend_box.text_frame
    trend_tf.word_wrap = True
    trend_tf.clear()

    obs_p = trend_tf.add_paragraph()
    obs_p.text = observation_title
    obs_p.font.size = Pt(14)
    obs_p.font.bold = True
    obs_p.font.color.rgb = RGBColor(255, 255, 255)

    add_bullets_with_bold(trend_tf, trend_texts, 10)

    rec_box = slide.shapes.add_textbox(
        left_width + Inches(0.2),
        Inches(2.8),
        right_width - Inches(0.4),
        Inches(2.8)
    )
    rec_tf = rec_box.text_frame
    rec_tf.word_wrap = True
    rec_tf.clear()

    rec_p = rec_tf.add_paragraph()
    rec_p.text = recommendation_title
    rec_p.font.size = Pt(14)
    rec_p.font.bold = True
    rec_p.font.color.rgb = RGBColor(255, 255, 255)

    add_bullets_with_bold(rec_tf, recommendations, 10)

    if position is not None:
        insert_slide_at(prs, slide, position)

def add_consolidated_slide(prs, trend_texts, recommendations, consolidated_title, observation_title, recommendation_title, position=None):
    blank_layout = next((layout for layout in prs.slide_layouts if layout.name.lower() == "blank"), prs.slide_layouts[6])
    slide = prs.slides.add_slide(blank_layout)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    background_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, slide_width, slide_height
    )
    background_shape.fill.solid()
    background_shape.fill.fore_color.rgb = RGBColor(255, 249, 237)
    background_shape.line.fill.background()
    slide.shapes._spTree.insert(2, background_shape._element)

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), slide_width - Inches(0.8), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = consolidated_title
    title_tf.paragraphs[0].font.size = Pt(20)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(30, 30, 50)

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), slide_width - Inches(1), slide_height - Inches(1.5))
    content_tf = content_box.text_frame
    content_tf.word_wrap = True

    p = content_tf.add_paragraph()
    p.text = observation_title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(50, 30, 30)

    add_bullets_with_bold(content_tf, trend_texts, 10, RGBColor(0, 0, 0))

    p = content_tf.add_paragraph()
    p.text = recommendation_title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(50, 30, 30)

    add_bullets_with_bold(content_tf, recommendations, 10, RGBColor(0, 0, 0))

    if position is not None:
        insert_slide_at(prs, slide, position)

def generate_ppt_report(summary, graph_insights, output_path, template_name=None):
    if template_name is None:
        template_name = selected_template
    
    template_path = os.path.join(os.path.dirname(__file__), 'templates', template_name)
    if not os.path.exists(template_path):
        template_path = os.path.join(os.path.dirname(__file__), 'templates', 'SMBC.pptx')
    
    prs = Presentation(template_path)

    first_slide = prs.slides[0]
    last_slide = prs.slides[-1]

    slides = list(prs.slides._sldIdLst)
    for i in range(len(slides) - 2, 0, -1):
        rId = slides[i].rId
        prs.part.drop_rel(rId)
        del slides[i]

    add_consolidated_slide(
        prs,
        trend_texts=summary['trend'],
        recommendations=summary['recommendation'],
        consolidated_title="Executive Summary",
        observation_title="Trend Observations:",
        recommendation_title="Recommendations:",
        position=1
    )

    for insight in graph_insights:
        filename = os.path.basename(insight["image_path"])
        name_without_ext = os.path.splitext(filename)[0]
        add_trend_slide(
            prs,
            image_path=insight["image_path"],
            trend_texts=insight["trend"],
            recommendations=insight["recommendation"],
            trend_title=insight["title"],
            observation_title="Trend Observations:",
            recommendation_title="Recommendations:",
            position=2
        )

    prs.save(output_path)

# API endpoint functions
async def get_templates():
    """Get available templates with preview info"""
    try:
        templates = get_available_templates()
        return {"templates": templates, "selected": selected_template}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to get templates: {str(e)}")

async def select_template(template_data: dict):
    """Select a template for report generation"""
    global selected_template
    
    template_name = template_data.get("template")
    if not template_name:
        raise HTTPException(status_code=400, detail="Template name required")
    
    templates_dir = os.path.join(os.path.dirname(__file__), 'templates')
    template_path = os.path.join(templates_dir, template_name)
    
    if not os.path.exists(template_path):
        raise HTTPException(status_code=404, detail="Template not found")
    
    selected_template = template_name
    return {"message": f"Template selected: {template_name}", "selected": selected_template}

async def upload_images(files: List[UploadFile]):
    global temp_dir
    
    if temp_dir:
        shutil.rmtree(temp_dir, ignore_errors=True)
    
    temp_dir = tempfile.mkdtemp()
    uploaded_files = []
    
    for file in files:
        if not file.filename.lower().endswith(('.png', '.jpg', '.jpeg')):
            raise HTTPException(status_code=400, detail=f"Invalid file type: {file.filename}")
        
        file_path = os.path.join(temp_dir, file.filename)
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        uploaded_files.append(file_path)
    
    return {"message": f"Uploaded {len(uploaded_files)} files", "files": uploaded_files}

async def configure_cropping(config: dict):
    global crop_config
    crop_config.update(config)
    return {"message": "Cropping configuration updated", "config": crop_config}

async def analyze():
    global analysis_results, temp_dir
    
    if not temp_dir or not os.path.exists(temp_dir):
        raise HTTPException(status_code=400, detail="No images uploaded")
    
    if not os.getenv("OPENAI_API_KEY"):
        raise HTTPException(status_code=500, detail="OpenAI API key not set")
    
    try:
        print(f"Analyzing images in: {temp_dir}")
        files = [f for f in os.listdir(temp_dir) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
        print(f"Found {len(files)} image files: {files}")
        
        if not files:
            raise HTTPException(status_code=400, detail="No valid image files found")
        
        # Crop images into blocks if enabled
        if crop_config.get("enabled", False):
            try:
                print("Cropping images into blocks...")
                cropped_files = process_uploaded_images(temp_dir, crop_config["rows"], crop_config["cols"])
                print(f"Created {len(cropped_files)} image blocks")
            except Exception as crop_error:
                print(f"Cropping failed: {crop_error}, continuing with original images")
        else:
            print("Cropping disabled, analyzing original images only")
        
        exec_summary, graph_insights = analyze_all_graphs(temp_dir)
        print(f"Analysis complete. Executive summary: {exec_summary}")
        print(f"Graph insights count: {len(graph_insights)}")
        
        analysis_results = {
            "executive_summary": exec_summary,
            "graph_insights": [
                {
                    "title": insight["title"],
                    "trend": insight["trend"],
                    "recommendation": insight["recommendation"],
                    "filename": os.path.basename(insight["image_path"]),
                    "image_url": f"/get-image/{os.path.basename(insight['image_path'])}"
                }
                for insight in graph_insights
            ]
        }
        
        return analysis_results
    except Exception as e:
        print(f"Analysis error: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Analysis failed: {str(e)}")

async def download_report(request=None):
    global analysis_results, temp_dir
    
    if not analysis_results:
        raise HTTPException(status_code=400, detail="No analysis results available")
    
    try:
        print("Starting report generation...")
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        report_filename = f"SMBC_Report_{timestamp}.pptx"
        report_path = os.path.join(tempfile.gettempdir(), report_filename)
        
        print(f"Report will be saved to: {report_path}")
        
        graph_insights_full = []
        for insight in analysis_results["graph_insights"]:
            full_path = None
            for root, dirs, files in os.walk(temp_dir):
                if insight["filename"] in files:
                    full_path = os.path.join(root, insight["filename"])
                    break
            
            if not full_path:
                full_path = os.path.join(temp_dir, insight["filename"])
            
            print(f"Image path: {full_path}, exists: {os.path.exists(full_path)}")
            if os.path.exists(full_path):
                graph_insights_full.append({
                    "image_path": full_path,
                    "title": insight["title"],
                    "trend": insight["trend"],
                    "recommendation": insight["recommendation"]
                })
        
        print("Calling generate_ppt_report...")
        if graph_insights_full:
            generate_ppt_report(
                analysis_results["executive_summary"],
                graph_insights_full,
                report_path,
                selected_template
            )
        else:
            raise Exception("No valid images found for report generation")
        
        print(f"Report generated, file exists: {os.path.exists(report_path)}")
        if os.path.exists(report_path):
            print(f"File size: {os.path.getsize(report_path)} bytes")
        
        return FileResponse(
            report_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=report_filename,
            headers={"Content-Disposition": f"attachment; filename={report_filename}"}
        )
    except Exception as e:
        print(f"Report generation error: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Report generation failed: {str(e)}")

async def get_image(filename: str):
    global temp_dir
    
    if not temp_dir:
        raise HTTPException(status_code=404, detail="No images available")
    
    # Search for the image in temp_dir and subdirectories
    for root, dirs, files in os.walk(temp_dir):
        if filename in files:
            image_path = os.path.join(root, filename)
            if os.path.exists(image_path):
                return FileResponse(image_path, media_type="image/jpeg")
    
    raise HTTPException(status_code=404, detail="Image not found")